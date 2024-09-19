import os
from pptx import Presentation
from pydub import AudioSegment
from tkinter import Tk
from tkinter.filedialog import askopenfilename, asksaveasfilename
from gtts import gTTS
import tempfile

def generate_slide_number_audio(slide_number):
    """Generate 'Slide X' audio using gTTS"""
    text = f"Slide {slide_number}"
    tts = gTTS(text)
    
    # Create a temporary audio file for the slide number
    with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as temp_audio:
        temp_audio_path = temp_audio.name
        tts.save(temp_audio_path)
    
    # Load the generated audio file as an AudioSegment
    slide_audio = AudioSegment.from_file(temp_audio_path)
    
    # Remove the temporary file
    os.remove(temp_audio_path)
    
    return slide_audio

def extract_audio_from_pptx(pptx_path):
    # Load the presentation
    prs = Presentation(pptx_path)
    
    # A list to store the audio segments
    audio_segments = []

    for i, slide in enumerate(prs.slides, start=1):
        # Add the 'Slide X' audio before each slide
        slide_number_audio = generate_slide_number_audio(i)
        audio_segments.append(slide_number_audio)
        
        for shape in slide.shapes:
            if shape.has_media and shape.media:
                # Save the embedded media to a temp file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.wav') as temp_audio:
                    temp_audio.write(shape.media.blob)
                    temp_audio_path = temp_audio.name
                    
                # Load the audio segment
                audio = AudioSegment.from_file(temp_audio_path)
                audio_segments.append(audio)
                
                # Remove the temporary file after loading the audio
                os.remove(temp_audio_path)
    
    # Concatenate all audio segments into one long segment
    if audio_segments:
        final_audio = sum(audio_segments)
        return final_audio
    else:
        return None

def main():
    # Tkinter setup
    root = Tk()
    root.withdraw()  # Hide the root window

    # Ask the user to select the PPTX file
    pptx_path = askopenfilename(
        title="Select the PowerPoint file",
        filetypes=[("PowerPoint files", "*.pptx")]
    )

    if not pptx_path:
        print("No file selected, exiting.")
        return

    # Extract the audio from the selected PPTX file
    final_audio = extract_audio_from_pptx(pptx_path)

    if final_audio:
        # Ask where to save the resulting audio
        save_path = asksaveasfilename(
            title="Save the combined audio",
            defaultextension=".wav",
            filetypes=[("WAV files", "*.wav")]
        )

        if save_path:
            final_audio.export(save_path, format="wav")
            print(f"Audio saved to {save_path}")
        else:
            print("No save location provided.")
    else:
        print("No audio found in the PowerPoint file.")

if __name__ == "__main__":
    main()
